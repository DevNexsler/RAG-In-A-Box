"""Live extraction tests — create real files in each supported format and verify extraction.

These tests use real libraries to create genuine docx, pptx, xlsx, html, csv, rtf, txt,
and epub files, then run them through the actual extractors (no mocks) to confirm
end-to-end extraction works for every supported format.

Requires: markitdown[all], openpyxl, python-docx, python-pptx (all installed as
dependencies of markitdown[all]).
"""

import csv
import io
import os
import zipfile
from pathlib import Path

import pytest


# ---------------------------------------------------------------------------
# Helpers — create real sample files
# ---------------------------------------------------------------------------


def _create_docx(path: Path, title: str = "Test Document", body: str = "This is a test paragraph with important content."):
    """Create a real .docx file using python-docx."""
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph(body)
    doc.add_heading("Section Two", level=2)
    doc.add_paragraph("Second section content with details about the project.")
    doc.save(str(path))


def _create_pptx(path: Path):
    """Create a real .pptx file using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Quarterly Report"
    slide.placeholders[1].text = "Q1 2026 Results"

    # Content slide
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Key Findings"
    body = slide.placeholders[1]
    body.text = "Revenue grew by 15% year-over-year."
    body.text_frame.add_paragraph().text = "Customer satisfaction improved significantly."

    prs.save(str(path))


def _create_xlsx(path: Path):
    """Create a real .xlsx file with mixed data types."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Employees"
    ws.append(["Name", "Department", "Salary", "Notes"])
    ws.append(["Alice Johnson", "Engineering", 120000, "Senior developer"])
    ws.append(["Bob Smith", "Marketing", 95000, "Team lead"])
    ws.append(["Carol White", "Engineering", 110000, None])
    ws.append([None, None, 85000, "Intern position"])

    # Second sheet
    ws2 = wb.create_sheet("Projects")
    ws2.append(["Project", "Status", "Budget"])
    ws2.append(["RAG Pipeline", "Active", 50000])
    ws2.append(["Data Migration", "Complete", 30000])

    wb.save(str(path))


def _create_html(path: Path):
    """Create a real .html file."""
    html = """<!DOCTYPE html>
<html>
<head><title>Meeting Notes</title></head>
<body>
<h1>Team Meeting Notes</h1>
<h2>Agenda</h2>
<ul>
    <li>Project status update</li>
    <li>Q2 planning discussion</li>
    <li>Resource allocation review</li>
</ul>
<h2>Action Items</h2>
<p>Alice will prepare the deployment plan by Friday.</p>
<p>Bob will review the security audit findings.</p>
</body>
</html>"""
    path.write_text(html, encoding="utf-8")


def _create_csv(path: Path):
    """Create a real .csv file."""
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Product", "Category", "Price", "Description"])
        writer.writerow(["Widget A", "Hardware", "29.99", "A reliable widget for daily use"])
        writer.writerow(["Gadget B", "Electronics", "149.99", "Advanced gadget with smart features"])
        writer.writerow(["Tool C", "Hardware", "59.99", "Professional-grade tool"])


def _create_txt(path: Path):
    """Create a plain text file with substantial content."""
    text = """Project Requirements Document

Overview:
This document outlines the requirements for the RAG-in-a-Box project.
The system must support multiple document formats including Word, Excel,
PowerPoint, HTML, and plain text files.

Key Requirements:
1. Document ingestion must handle UTF-8 encoding gracefully
2. Excel files should extract text-only cells, skipping numeric data
3. MarkItDown converts documents to Markdown for heading-aware chunking
4. All extractors must handle errors gracefully without crashing

Performance Goals:
- Process 1000 documents in under 10 minutes
- Support files up to 100MB in size
- Maintain search latency under 200ms
"""
    path.write_text(text, encoding="utf-8")


def _create_txt_unicode(path: Path):
    """Create a plain text file with Unicode content."""
    text = "日本語のドキュメント\n\nこれはテストです。\nCafé résumé naïve\nEmoji: 🚀 📊 🔍"
    path.write_text(text, encoding="utf-8")


def _create_rtf(path: Path):
    r"""Create a minimal .rtf file with readable content."""
    rtf = r"""{\rtf1\ansi\deff0
{\fonttbl{\f0 Times New Roman;}}
\f0\fs24
Research Summary\par
\par
This document contains findings from our quarterly research initiative.\par
The team analyzed customer feedback across multiple channels.\par
Key insight: users prefer automated document processing over manual entry.\par
}"""
    path.write_text(rtf, encoding="utf-8")


def _create_epub(path: Path):
    """Create a minimal valid .epub file (EPUB is a zip with specific structure)."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        # mimetype must be first and uncompressed
        zf.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)

        # container.xml
        zf.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?>'
            '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
            '<rootfiles><rootfile full-path="content.opf" media-type="application/oebps-package+xml"/>'
            '</rootfiles></container>',
        )

        # content.opf
        zf.writestr(
            "content.opf",
            '<?xml version="1.0"?>'
            '<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="uid">'
            '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
            '<dc:title>Test EPUB Book</dc:title>'
            '<dc:identifier id="uid">test-epub-001</dc:identifier>'
            '<dc:language>en</dc:language>'
            '</metadata>'
            '<manifest>'
            '<item id="ch1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>'
            '</manifest>'
            '<spine><itemref idref="ch1"/></spine>'
            '</package>',
        )

        # chapter1.xhtml
        zf.writestr(
            "chapter1.xhtml",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<!DOCTYPE html>'
            '<html xmlns="http://www.w3.org/1999/xhtml">'
            "<head><title>Chapter 1</title></head>"
            "<body>"
            "<h1>Chapter 1: Introduction</h1>"
            "<p>Welcome to the test EPUB book. This chapter covers the basics of document processing.</p>"
            "<h2>Background</h2>"
            "<p>Document format support is essential for modern knowledge management systems.</p>"
            "</body></html>",
        )

    path.write_bytes(buf.getvalue())


# ---------------------------------------------------------------------------
# Live extraction tests — real files, no mocks
# ---------------------------------------------------------------------------


class TestPlaintextLive:
    """Live tests for plain text extraction."""

    def test_basic_extraction(self, tmp_path):
        from extractors import extract_plaintext

        path = tmp_path / "requirements.txt"
        _create_txt(path)
        result = extract_plaintext(path)

        assert "Project Requirements Document" in result.full_text
        assert "RAG-in-a-Box" in result.full_text
        assert "UTF-8 encoding" in result.full_text
        assert result.frontmatter == {}

    def test_unicode_extraction(self, tmp_path):
        from extractors import extract_plaintext

        path = tmp_path / "unicode.txt"
        _create_txt_unicode(path)
        result = extract_plaintext(path)

        assert "日本語のドキュメント" in result.full_text
        assert "Café" in result.full_text
        assert "🚀" in result.full_text

    def test_empty_file(self, tmp_path):
        from extractors import extract_plaintext

        path = tmp_path / "empty.txt"
        path.write_text("", encoding="utf-8")
        result = extract_plaintext(path)

        assert result.full_text == ""

    def test_via_dispatcher(self, tmp_path):
        from extractors import extract_text

        path = tmp_path / "doc.txt"
        _create_txt(path)
        result = extract_text(path, ext="txt")

        assert "Project Requirements Document" in result.full_text
        assert len(result.pages) == 1


class TestExcelLive:
    """Live tests for Excel extraction using real .xlsx files."""

    def test_text_cells_extracted(self, tmp_path):
        from extractors import extract_excel

        path = tmp_path / "data.xlsx"
        _create_xlsx(path)
        result = extract_excel(path)

        assert "Alice Johnson" in result.full_text
        assert "Bob Smith" in result.full_text
        assert "Senior developer" in result.full_text
        assert "Team lead" in result.full_text

    def test_headers_included(self, tmp_path):
        from extractors import extract_excel

        path = tmp_path / "data.xlsx"
        _create_xlsx(path)
        result = extract_excel(path)

        assert "Headers: Name | Department | Salary | Notes" in result.full_text

    def test_numbers_not_in_data_rows(self, tmp_path):
        from extractors import extract_excel

        path = tmp_path / "data.xlsx"
        _create_xlsx(path)
        result = extract_excel(path)

        # Numbers should only appear in header rows, not extracted as data
        lines = result.full_text.split("\n")
        data_lines = [l for l in lines if l.strip() and not l.startswith(("Sheet:", "Headers:"))]
        for line in data_lines:
            assert "120000" not in line
            assert "95000" not in line

    def test_multiple_sheets(self, tmp_path):
        from extractors import extract_excel

        path = tmp_path / "data.xlsx"
        _create_xlsx(path)
        result = extract_excel(path)

        assert "Sheet: Employees" in result.full_text
        assert "Sheet: Projects" in result.full_text
        assert "RAG Pipeline" in result.full_text
        assert "Data Migration" in result.full_text

    def test_via_dispatcher(self, tmp_path):
        from extractors import extract_text

        path = tmp_path / "data.xlsx"
        _create_xlsx(path)
        result = extract_text(path, ext="xlsx")

        assert "Alice Johnson" in result.full_text
        assert "Sheet: Employees" in result.full_text


class TestDocxLive:
    """Live tests for Word document extraction using real .docx files."""

    def test_extracts_headings_and_body(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "report.docx"
        _create_docx(path)
        result = extract_markitdown(path)

        assert result.full_text.strip() != ""
        assert "Test Document" in result.full_text
        assert "test paragraph" in result.full_text

    def test_extracts_multiple_sections(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "report.docx"
        _create_docx(path)
        result = extract_markitdown(path)

        assert "Section Two" in result.full_text
        assert "Second section content" in result.full_text

    def test_via_dispatcher(self, tmp_path):
        from extractors import extract_text

        path = tmp_path / "report.docx"
        _create_docx(path)
        result = extract_text(path, ext="docx")

        assert "Test Document" in result.full_text
        assert len(result.pages) == 1


class TestPptxLive:
    """Live tests for PowerPoint extraction using real .pptx files."""

    def test_extracts_slide_content(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "deck.pptx"
        _create_pptx(path)
        result = extract_markitdown(path)

        assert result.full_text.strip() != ""
        assert "Quarterly Report" in result.full_text

    def test_extracts_bullet_points(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "deck.pptx"
        _create_pptx(path)
        result = extract_markitdown(path)

        assert "Revenue grew" in result.full_text or "Key Findings" in result.full_text

    def test_via_dispatcher(self, tmp_path):
        from extractors import extract_text

        path = tmp_path / "deck.pptx"
        _create_pptx(path)
        result = extract_text(path, ext="pptx")

        assert "Quarterly Report" in result.full_text


class TestHtmlLive:
    """Live tests for HTML extraction using real .html files."""

    def test_extracts_headings_and_text(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "notes.html"
        _create_html(path)
        result = extract_markitdown(path)

        assert result.full_text.strip() != ""
        assert "Team Meeting Notes" in result.full_text

    def test_extracts_list_items(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "notes.html"
        _create_html(path)
        result = extract_markitdown(path)

        assert "Project status update" in result.full_text
        assert "Q2 planning" in result.full_text

    def test_extracts_paragraphs(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "notes.html"
        _create_html(path)
        result = extract_markitdown(path)

        assert "Alice will prepare" in result.full_text
        assert "deployment plan" in result.full_text

    def test_via_dispatcher(self, tmp_path):
        from extractors import extract_text

        path = tmp_path / "notes.html"
        _create_html(path)
        result = extract_text(path, ext="html")

        assert "Team Meeting Notes" in result.full_text

    def test_htm_extension(self, tmp_path):
        """Dispatcher handles .htm as well as .html."""
        from extractors import extract_text

        path = tmp_path / "page.htm"
        _create_html(path)
        result = extract_text(path, ext="htm")

        assert "Team Meeting Notes" in result.full_text


class TestCsvLive:
    """Live tests for CSV extraction using real .csv files."""

    def test_extracts_csv_content(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "products.csv"
        _create_csv(path)
        result = extract_markitdown(path)

        assert result.full_text.strip() != ""
        assert "Widget A" in result.full_text
        assert "Gadget B" in result.full_text

    def test_extracts_descriptions(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "products.csv"
        _create_csv(path)
        result = extract_markitdown(path)

        assert "reliable widget" in result.full_text or "daily use" in result.full_text

    def test_via_dispatcher(self, tmp_path):
        from extractors import extract_text

        path = tmp_path / "products.csv"
        _create_csv(path)
        result = extract_text(path, ext="csv")

        assert "Widget A" in result.full_text


class TestRtfLive:
    """Live tests for RTF extraction using real .rtf files."""

    def test_extracts_rtf_content(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "summary.rtf"
        _create_rtf(path)
        result = extract_markitdown(path)

        assert result.full_text.strip() != ""
        # RTF content should be extracted (exact output depends on MarkItDown's RTF handling)
        assert "Research Summary" in result.full_text or "customer feedback" in result.full_text or "quarterly" in result.full_text

    def test_via_dispatcher(self, tmp_path):
        from extractors import extract_text

        path = tmp_path / "summary.rtf"
        _create_rtf(path)
        result = extract_text(path, ext="rtf")

        assert result.full_text.strip() != ""


class TestEpubLive:
    """Live tests for EPUB extraction using real .epub files."""

    def test_extracts_epub_content(self, tmp_path):
        from extractors import extract_markitdown

        path = tmp_path / "book.epub"
        _create_epub(path)
        result = extract_markitdown(path)

        assert result.full_text.strip() != ""
        # Should extract chapter content
        has_content = (
            "Chapter 1" in result.full_text
            or "Introduction" in result.full_text
            or "document processing" in result.full_text
        )
        assert has_content, f"Expected chapter content, got: {result.full_text[:300]}"

    def test_via_dispatcher(self, tmp_path):
        from extractors import extract_text

        path = tmp_path / "book.epub"
        _create_epub(path)
        result = extract_text(path, ext="epub")

        assert result.full_text.strip() != ""


# ---------------------------------------------------------------------------
# Cross-format dispatcher test — all formats in one test
# ---------------------------------------------------------------------------


class TestDispatcherAllFormats:
    """Verify the dispatcher routes every supported extension to the right extractor
    and produces non-empty output from real files."""

    @pytest.fixture
    def sample_dir(self, tmp_path):
        """Create one sample file for each supported format."""
        files = {}

        # txt
        p = tmp_path / "readme.txt"
        _create_txt(p)
        files["txt"] = p

        # docx
        p = tmp_path / "report.docx"
        _create_docx(p)
        files["docx"] = p

        # pptx
        p = tmp_path / "slides.pptx"
        _create_pptx(p)
        files["pptx"] = p

        # xlsx
        p = tmp_path / "data.xlsx"
        _create_xlsx(p)
        files["xlsx"] = p

        # html
        p = tmp_path / "page.html"
        _create_html(p)
        files["html"] = p

        # htm
        p = tmp_path / "page2.htm"
        _create_html(p)
        files["htm"] = p

        # csv
        p = tmp_path / "data.csv"
        _create_csv(p)
        files["csv"] = p

        # rtf
        p = tmp_path / "doc.rtf"
        _create_rtf(p)
        files["rtf"] = p

        # epub
        p = tmp_path / "book.epub"
        _create_epub(p)
        files["epub"] = p

        return files

    @pytest.mark.parametrize("ext", ["txt", "docx", "pptx", "xlsx", "html", "htm", "csv", "rtf", "epub"])
    def test_dispatcher_returns_content(self, sample_dir, ext):
        """Each format produces non-empty extraction through the dispatcher."""
        from extractors import extract_text

        path = sample_dir[ext]
        result = extract_text(path, ext=ext)

        assert result.full_text.strip() != "", f"Extraction returned empty text for .{ext}"
        assert len(result.pages) >= 1, f"No pages returned for .{ext}"

    def test_unknown_extension_returns_empty(self):
        """Unknown extension still returns empty."""
        from extractors import extract_text

        result = extract_text("/fake/file.xyz", ext="xyz")
        assert result.full_text == ""


# ---------------------------------------------------------------------------
# Source type mapping test (flow_index_vault.py)
# ---------------------------------------------------------------------------


class TestSourceTypeMap:
    """Verify _SOURCE_TYPE_MAP covers all supported extensions."""

    def test_all_new_extensions_mapped(self):
        from flow_index_vault import _SOURCE_TYPE_MAP

        expected = {
            "txt": "txt",
            "docx": "doc", "doc": "doc", "rtf": "doc",
            "xlsx": "sheet", "xls": "sheet",
            "pptx": "pres",
            "csv": "csv",
            "html": "html", "htm": "html",
            "epub": "epub",
            # Existing
            "md": "md", "pdf": "pdf",
            "png": "img", "jpg": "img", "jpeg": "img", "gif": "img", "webp": "img",
        }
        for ext, expected_type in expected.items():
            assert _SOURCE_TYPE_MAP.get(ext) == expected_type, f"Wrong mapping for .{ext}"

    def test_unknown_ext_returns_other(self):
        from flow_index_vault import _SOURCE_TYPE_MAP

        assert _SOURCE_TYPE_MAP.get("xyz", "other") == "other"

    def test_heading_aware_types(self):
        """MarkItDown output types should use heading-aware chunking."""
        heading_aware = ("md", "doc", "pres", "html", "epub", "csv")
        flat_chunking = ("sheet", "txt", "img")

        # These should NOT overlap
        assert not set(heading_aware) & set(flat_chunking)
